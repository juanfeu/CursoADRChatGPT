import random
from openai import OpenAI
from pptx import Presentation
from docx import Document
from pptx.util import Inches

# Initialize the OpenAI client with your API key
client = OpenAI(
    api_key="",
)

# List of climate change-related words
climate_words = ["Ecosystem", "Biodiversity", "Sustainability", "Greenhouse", "Carbon", "Renewable", "Emissions",
                 "Conservation", "Pollution", "Recycle"]


def get_summary(text):
    """ Use the ChatGPT API to get a summary of the provided text """
    print("Debug: Summarizing text:", text[:50])  # Show only the first 50 characters for brevity
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": f"Summarize this text: {text}"}],
            temperature=0.7,  # Adjust creativity of the responses.
            max_tokens=50,  # Limit the length of the generated response.
            top_p=1.0,  # Control the diversity via nucleus sampling.
            frequency_penalty=0.1,  # Reduce repetition of words and phrases.
            presence_penalty=0.0  # Encourage the introduction of new concepts.
        )
        if response.choices:
            summarized_text = response.choices[0].message.content
            print("Debug: Summarized text:", summarized_text[:50])
            return summarized_text
        else:
            print("Debug: No choices available in the response.")
            return "Summary could not be generated."
    except Exception as e:
        print("Error in summarization:", str(e))
        return "Summary could not be generated."


def doc_to_pptx(doc_path, pptx_path):
    doc = Document(doc_path)
    ppt = Presentation()

    print("Debug: Starting document processing.")
    for para in doc.paragraphs:
        if para.text.strip() and len(para.text.strip()) > 40:  # Ensure paragraph has content and is significant
            print("Debug: Processing paragraph with text:", para.text[:50])
            summarized_text = get_summary(para.text)
            if summarized_text != "Summary could not be generated.":
                slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and content layout

                # Choose a random climate change-related word as the title
                title_text = random.choice(climate_words)
                title = slide.shapes.title
                title.text = title_text
                print("Debug: Title added to slide:", title_text)

                # Add the full summary as the main content
                if len(slide.shapes.placeholders) > 1:
                    content = slide.shapes.placeholders[1]
                    content.text = summarized_text
                    print("Debug: Content added to slide:", summarized_text[:50])

    ppt.save(pptx_path)
    print("Debug: PowerPoint saved to", pptx_path)


# Paths to the Word document and the PowerPoint file to be created
doc_path = r"C:\openai\word\your_word_document.docx"
pptx_path = r"C:\openai\powerpoint\your_powerpoint.pptx"

doc_to_pptx(doc_path, pptx_path)









