import random
import requests
from io import BytesIO
from PIL import Image
from openai import OpenAI
from pptx import Presentation
from docx import Document
from pptx.util import Inches

# Initialize the OpenAI client with your API key
client = OpenAI(
    api_key="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",
)

# List of climate change-related words
climate_words = ["Ecosistema", "Biodiversidad", "Sostenibilidad", "Greenhouse", "Carbono", "Renovable", "Emisiones", "Conservación", "Contaminación", "Reciclaje"]

def get_image(description):
    """ Generate an image based on a description using the OpenAI API """
    url = "https://api.openai.com/v1/images/generations"
    headers = {
        "Authorization": f"Bearer {client.api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": "dall-e-2",
        "prompt": description,
        "n": 1,
        "size": "256x256"  # Smallest size typically offered; resize later to 75x75
    }
    response = requests.post(url, json=payload, headers=headers)
    if response.status_code == 200:
        image_data = response.json()
        image_url = image_data['data'][0]['url']
        image_response = requests.get(image_url)
        if image_response.status_code == 200:
            image = Image.open(BytesIO(image_response.content))
            image = image.resize((50, 50), Image.Resampling.LANCZOS)  # Updated to use LANCZOS for high quality
            img_byte_arr = BytesIO()
            image.save(img_byte_arr, format='PNG')
            return img_byte_arr
        else:
            print("Error al descargar la imagen")
    else:
        print(f"Error al crear la imagen: {response.status_code} - {response.text}")
    return None

def get_summary(text):
    """ ChatGPT te genera un resumen a partir del texto previsto """
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": f"Resume este texto sin utilizar más de 15 palabras: {text}"}],
            temperature=0.7,
            max_tokens=100,
            top_p=1.0,
            frequency_penalty=0.1,
            presence_penalty=0.0
        )
        if response.choices:
            summarized_text = response.choices[0].message.content
            return summarized_text
        else:
            return "No se pudo general el resumen."
    except Exception as e:
        print("Error al resumir:", str(e))
        return "No se pudo generar al resumen."

def doc_to_pptx(doc_path, pptx_path):
    doc = Document(doc_path)
    ppt = Presentation()
    first_slide = True

    for para in doc.paragraphs:
        if para.text.strip() and len(para.text.strip()) > 40:
            summarized_text = get_summary(para.text)
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and content layout

            if first_slide and summarized_text != "No se pudo crear el resumen":
                image_stream = get_image("Crea una imagen que tenga que ver con el medio ambiente y que no tenga texto")
                if image_stream:
                    slide.shapes.add_picture(image_stream, Inches(1), Inches(1), width=Inches(1.5), height=Inches(1.5))  # Adjust size as needed
                    first_slide = False

            if summarized_text != "No se pudo generar el resumen.":
                title_text = random.choice(climate_words)
                slide.shapes.title.text = title_text
                content = slide.shapes.placeholders[1]
                content.text = summarized_text

    ppt.save(pptx_path)
    print("PowerPoint saved to", pptx_path)

# Paths to the Word document and the PowerPoint file to be created
doc_path = r"C:\openai\word\your_word_document.docx"
pptx_path = r"C:\openai\powerpoint\your_powerpoint.pptx"

doc_to_pptx(doc_path, pptx_path)











